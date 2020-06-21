from selenium import webdriver
from selenium.webdriver.common.keys import Keys
from selenium.webdriver.firefox.options import Options
from pptx import Presentation
from pptx.util import Inches, Pt
import time
def login(browser,username,passwd):
    input = browser.find_element_by_name("user")
    input.send_keys(username)
    input = browser.find_element_by_name("password")
    input.send_keys(passwd)
    time.sleep(3)
    input.send_keys(Keys.ENTER)
if __name__ == "__main__":
    prs = Presentation()
    prs.slide_height = 6858000
    prs.slide_width = 12192000
    username="" #登入grafana帳號
    passwd=""   #登入grafana密碼
    options = Options()
    browser = webdriver.Firefox(executable_path="/Users/Desktop/geckodriver",options=options)
    browser.maximize_window()
    screenshotlist = {"path":"path2","path3":"path4"}
    time.sleep(2)
    furl = "https://abc.asc.com"
    browser.get(furl)
    login(browser,username,passwd)
    time.sleep(3)
    dest = "https://aaa.bbb.com"
    for i in screenshotlist:
        bullet_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        left = width = height = Inches(1)
        textbox = slide.shapes.add_textbox(left, Inches(0), width, height)
        tf = textbox.text_frame
        p = tf.add_paragraph()
        p.text = i
        p.font.size = Pt(40)
        url = dest+i+"&var-node="+screenshotlist[i]+"&var-port=9100"
        browser.get(url)
        time.sleep(7)
        ele = browser.find_element_by_class_name('panel-height-helper')
        ele.screenshot(i+".png")
        time.sleep(3)
        slide.shapes.add_picture(i+".png", Inches(0.5), Inches(1), Inches(12), Inches(6))
        print("done")
    prs.save("grafana.pptx")
    browser.quit()
